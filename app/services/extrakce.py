"""services/extrakce.py — extrakce textu z různých formátů (PDF/Word/Excel/PPT/txt).
Převzato z aplikace Brain (osvědčené)."""
import io


PODPOROVANE = ("pdf", "docx", "xlsx", "xls", "pptx", "txt")


def lze_extrahovat(filename):
    return filename.rsplit(".", 1)[-1].lower() in PODPOROVANE if "." in filename else False


def extrahuj_text(file_bytes: bytes, filename: str) -> str:
    """Vytáhne text ze souboru podle přípony."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    try:
        if ext == "pdf":
            return _z_pdf(file_bytes)
        elif ext == "docx":
            return _z_docx(file_bytes)
        elif ext in ("xlsx", "xls"):
            return _z_excel(file_bytes)
        elif ext == "pptx":
            return _z_pptx(file_bytes)
        elif ext == "txt":
            return file_bytes.decode("utf-8", errors="ignore")
        return ""
    except Exception as e:
        return f"[Chyba při extrakci: {e}]"


def _z_pdf(data: bytes) -> str:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        return "\n\n".join((p.extract_text() or "") for p in reader.pages)
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return "\n\n".join(p.extract_text() or "" for p in reader.pages)
    except Exception as e:
        return f"[PDF nelze přečíst: {e}]"


def _z_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _z_excel(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    radky = []
    for sheet in wb.worksheets:
        radky.append(f"--- List: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None]
            if vals:
                radky.append("\t".join(vals))
    return "\n".join(radky)


def _z_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    texty = []
    for i, slide in enumerate(prs.slides, 1):
        st = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        if st:
            texty.append(f"--- Slide {i} ---\n" + "\n".join(st))
    return "\n\n".join(texty)
